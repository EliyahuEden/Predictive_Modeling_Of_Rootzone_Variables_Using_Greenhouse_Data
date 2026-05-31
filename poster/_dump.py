"""Dump all shapes on slide 1 of the Project 14 base file (id, name, geometry, text)."""
import os
from pptx import Presentation
from pptx.util import Emu

ROOT = r"C:\Users\User\Desktop\Predictive_Modeling_Of_Rootzone_Variables_Using_Greenhouse_Data"
DELIV = os.path.join(ROOT, "poster",
                     "Predictive Modeling of Root-Zone Variables Using Greenhouse Data - Project 14.pptx")

def cm(v):
    return None if v is None else round(Emu(v).cm, 2)

prs = Presentation(DELIV)
print(f"slide size: {round(prs.slide_width/360000,2)} x {round(prs.slide_height/360000,2)} cm")
slide = prs.slides[0]
for sh in slide.shapes:
    txt = ""
    if sh.has_text_frame:
        txt = " | ".join(p.text for p in sh.text_frame.paragraphs if p.text).strip()
    txt = (txt[:70] + "...") if len(txt) > 70 else txt
    print(f"#{sh.shape_id:<4} {sh.shape_type!s:<22} L={cm(sh.left)} T={cm(sh.top)} "
          f"W={cm(sh.width)} H={cm(sh.height)}  {txt}")
