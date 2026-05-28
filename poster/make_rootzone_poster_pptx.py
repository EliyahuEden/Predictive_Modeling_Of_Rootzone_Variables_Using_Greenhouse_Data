from io import BytesIO
from pathlib import Path
from zipfile import ZipFile

from PIL import Image, ImageChops
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Emu, Pt


DPI = 300
W_CM = 70
H_CM = 100
W = round(W_CM / 2.54 * DPI)
H = round(H_CM / 2.54 * DPI)

ROOT = Path(__file__).resolve().parents[1]
POSTER_DIR = ROOT / "poster"
PLOTS = ROOT / "plots"
OUT_PPTX = POSTER_DIR / "rootzone_poster_70x100cm_editable.pptx"

BLUE = "#1976D2"
TEAL = "#006B82"
DARK = "#0D2B38"
DEEP_TEAL = "#006276"
GREEN = "#75C856"
OLIVE = "#A6C946"
BORDER = "#9FD4E8"
TEXT_MUTED = "#4D6C7C"
LINE = "#A8D6E8"
WHITE = "#FFFFFF"


def emu(v):
    return Emu(round(v * 914400 / DPI))


def rgb(hex_color):
    return RGBColor.from_string(hex_color.strip("#"))


def line_pt(px):
    return Pt(px * 72 / DPI)


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)


def set_line(shape, color=None, width_px=3):
    if color is None:
        shape.line.fill.background()
        return
    shape.line.color.rgb = rgb(color)
    shape.line.width = line_pt(width_px)


def add_round(slide, x, y, w, h, fill, outline=BORDER, width_px=3):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(x), emu(y), emu(w), emu(h))
    set_fill(shp, fill)
    set_line(shp, outline, width_px)
    return shp


def add_rect(slide, x, y, w, h, fill, outline=None, width_px=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(x), emu(y), emu(w), emu(h))
    set_fill(shp, fill)
    set_line(shp, outline, width_px)
    return shp


def add_oval(slide, x, y, w, h, fill, outline=None, width_px=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, emu(x), emu(y), emu(w), emu(h))
    set_fill(shp, fill)
    set_line(shp, outline, width_px)
    return shp


def add_text(slide, x, y, w, h, text, size, color=DARK, bold=False, align="left", valign="top", font="Arial"):
    box = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(valign, MSO_ANCHOR.TOP)
    alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(align, PP_ALIGN.LEFT)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alignment
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def add_center_text(slide, x, y, w, h, text, size, color=DARK, bold=False):
    return add_text(slide, x, y, w, h, text, size, color, bold, "center", "middle")


def add_arrow(slide, x1, y1, x2, y2, color=LINE):
    if x2 <= x1:
        return
    h = 34
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, emu(x1), emu(y1 - h / 2), emu(x2 - x1), emu(h))
    set_fill(shp, color)
    set_line(shp, color, 1)
    return shp


def add_chip(slide, x, y, text, accent=BLUE, fill="#F6FDFF", w=None, h=120):
    if w is None:
        w = max(220, 34 * len(text) + 100)
    add_round(slide, x, y, w, h, fill, accent, 3)
    add_center_text(slide, x + 12, y + 6, w - 24, h - 12, text, 13.5, accent, True)
    return w


def add_metric(slide, x, y, w, h, value, label, accent):
    add_round(slide, x, y, w, h, "#F2FBFE", accent, 4)
    add_center_text(slide, x, y + 12, w, h * 0.56, value, 42, accent, True)
    add_center_text(slide, x + 10, y + h * 0.57, w - 20, h * 0.33, label, 14, TEXT_MUTED, True)


def trim_white(path):
    img = Image.open(path).convert("RGBA")
    rgb_img = img.convert("RGB")
    diff = ImageChops.difference(rgb_img, Image.new("RGB", rgb_img.size, "white"))
    bbox = diff.getbbox()
    if not bbox:
        return img
    pad = 18
    left = max(0, bbox[0] - pad)
    top = max(0, bbox[1] - pad)
    right = min(img.width, bbox[2] + pad)
    bottom = min(img.height, bbox[3] + pad)
    return img.crop((left, top, right, bottom))


def image_bytes(img):
    buf = BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def add_picture_contain(slide, img, x, y, w, h):
    if isinstance(img, (str, Path)):
        im = Image.open(img)
        source = Path(img)
    else:
        im = img
        source = image_bytes(img)
    scale = min(w / im.width, h / im.height)
    nw, nh = im.width * scale, im.height * scale
    return slide.shapes.add_picture(source, emu(x + (w - nw) / 2), emu(y + (h - nh) / 2), width=emu(nw), height=emu(nh))


def load_logos():
    ppt = POSTER_DIR / "hypotheticalocean.pptx"
    with ZipFile(ppt) as z:
        logos = []
        for name in z.namelist():
            if name.startswith("ppt/media/"):
                img = Image.open(BytesIO(z.read(name))).convert("RGBA")
                if img.size in [(217, 217), (900, 900)]:
                    logos.append(img)
    ruppin = next(i for i in logos if i.size == (217, 217))
    volcani = next(i for i in logos if i.size == (900, 900))
    return ruppin, volcani


def section(slide, x, y, w, h, n, title):
    add_round(slide, x, y, w, h, WHITE, BORDER, 6)
    add_oval(slide, x + 50, y + 70, 270, 270, BLUE)
    add_center_text(slide, x + 50, y + 70, 270, 270, str(n), 23, WHITE, True)
    add_text(slide, x + 360, y + 88, w - 450, 150, title, 36, TEAL, True)
    return x + 155, y + 345, w - 310, h - 435


def build_header(slide, ruppin, volcani):
    x, y, w, h = 230, 170, W - 460, 900
    add_round(slide, x, y, w, h, WHITE, "#4BA5BF", 7)
    add_picture_contain(slide, ruppin, x + 145, y + 95, 470, 470)
    add_picture_contain(slide, volcani, x + w - 615, y + 90, 470, 470)
    add_center_text(slide, x + 690, y + 80, w - 1380, 360, "Predictive Modeling of Rootzone Variables\nUsing Greenhouse Data", 44, TEAL, True)
    add_center_text(slide, x + 670, y + 485, w - 1340, 110, "Students: Eden Eliyahu, Oren Chaushu", 22, TEAL, True)
    add_center_text(slide, x + 670, y + 635, w - 1340, 110, "Mentors: Dr. Avner Priel - Ruppin,  Dr. Alaa Jamal - Volcani Institute", 18, DARK, True)


def build_intro(slide, x, y, w, h):
    cx, cy, cw, _ = section(slide, x, y, w, h, 1, "Introduction & Goal")
    add_text(slide, cx, cy, cw, 90, "Why rootzone monitoring needs a soft sensor.", 25, DARK, True)
    para = (
        "Rootzone pH and EC are key indicators of nutrient availability, salinity, and plant "
        "health in greenhouse cultivation. Manual samples are accurate, but they are sparse - "
        "between samples, acidity or salinity changes may go undetected."
    )
    add_text(slide, cx, cy + 210, cw, 330, para, 20, TEXT_MUTED)
    box_y = cy + 735
    bh = 790
    bw = (cw - 2 * 190) / 3
    labels = [
        ("Manual\nsample", "accurate -\nsparse", DEEP_TEAL, "#EFFAFF"),
        ("Blind\ninterval", "no observation", "#78910D", "#FBFEEB"),
        ("Rootzone\nsoft-sensor\nprediction", "dense -\nestimated", BLUE, "#F1F7FF"),
    ]
    for i, (title, sub, col, fill) in enumerate(labels):
        bx = cx + i * (bw + 190)
        add_round(slide, bx, box_y, bw, bh, fill, col, 4)
        add_oval(slide, bx + bw / 2 - 82, box_y + 60, 164, 164, fill, col, 16)
        add_center_text(slide, bx + 45, box_y + 280, bw - 90, 260, title, 20, col, True)
        add_center_text(slide, bx + 45, box_y + 535, bw - 90, 180, sub, 17, TEXT_MUTED)
        if i < 2:
            add_arrow(slide, bx + bw + 45, box_y + bh / 2, bx + bw + 125, box_y + bh / 2)
    goal_y = box_y + bh + 120
    add_round(slide, cx, goal_y, cw, 350, "#F4FCFE", BORDER, 4)
    add_oval(slide, cx + 80, goal_y + 135, 65, 65, "#A7E08F")
    add_oval(slide, cx + 96, goal_y + 151, 33, 33, GREEN)
    add_text(slide, cx + 210, goal_y + 75, cw - 300, 210, "Goal - Build a rootzone soft sensor that predicts future pH and EC between manual samples.", 20, TEAL, True)


def build_dataset(slide, x, y, w, h):
    cx, cy, cw, _ = section(slide, x, y, w, h, 2, "Dataset & Processing")
    add_text(slide, cx, cy, cw, 90, "Five greenhouse domains synchronized into one 10-minute master dataset.", 25, DARK, True)
    row_x, row_y = cx, cy + 285
    row_w, row_h, gap = int(cw * 0.58), 225, 52
    rows = [
        ("Climate conditions", "temperature - humidity - radiation - ...", BLUE),
        ("Irrigation & fertigation", "volume - fertilizer - acid & salt inputs", GREEN),
        ("Plant status", "canopy cover - crop age - growth stage", OLIVE),
        ("Rootzone measurements", "manual pH samples - manual EC samples", "#08798B"),
        ("Time & history", "timestamps - gaps - previous state - ...", "#1A5570"),
    ]
    for i, (lab, desc, col) in enumerate(rows):
        yy = row_y + i * (row_h + gap)
        add_round(slide, row_x, yy, row_w, row_h, WHITE, BORDER, 3)
        add_oval(slide, row_x + 60, yy + 82, 58, 58, col)
        add_text(slide, row_x + 155, yy + 62, row_w * 0.38, 80, lab, 17, TEAL, True)
        add_text(slide, row_x + row_w * 0.50, yy + 67, row_w * 0.48, 80, desc, 17, TEXT_MUTED)
    card_x, card_y = cx + int(cw * 0.63), row_y
    card_w, card_h = cw - (card_x - cx), 1330
    add_round(slide, card_x, card_y, card_w, card_h, DEEP_TEAL, "#034A58", 4)
    add_center_text(slide, card_x, card_y + 90, card_w, 320, "MASTER DATASET\n10-minute\ngreenhouse grid", 20, WHITE, True)
    add_round(slide, card_x, card_y + 540, card_w, card_h - 540, "#197388", "#197388", 1)
    stats = [("16,682", "ROWS"), ("10 min", "GRAIN"), ("109", "PH LABELS"), ("109", "EC LABELS"), ("20", "COLUMNS"), ("May-Sep", "2025")]
    sx = [card_x + 70, card_x + card_w * 0.52]
    sy = card_y + 625
    for i, (val, lab) in enumerate(stats):
        add_text(slide, sx[i % 2], sy + (i // 2) * 275, 430, 80, val, 18, WHITE, True)
        add_text(slide, sx[i % 2], sy + 100 + (i // 2) * 275, 430, 70, lab, 13, "#C6EFF7", True)
    add_arrow(slide, row_x + row_w + 35, row_y + 2 * (row_h + gap) + 112, card_x - 30, card_y + card_h / 2)
    density_y = card_y + card_h + 90
    add_round(slide, cx, density_y, cw, 535, "#EEF9FC", BORDER, 3)
    add_text(slide, cx + 90, density_y + 82, 1200, 60, "TELEMETRY DENSITY VS. LABEL SPARSITY", 13, "#688798", True)
    add_text(slide, cx + cw - 1110, density_y + 82, 850, 60, "LABELS: 109 PH - 109 EC", 13, TEAL, True, "right")
    bar_x, bar_y, bar_w, bar_h = cx + 90, density_y + 190, cw - 180, 78
    add_round(slide, bar_x, bar_y, bar_w, bar_h, "#BDEAF4", "#90CEDF", 2)
    for i in range(70):
        add_rect(slide, bar_x + 15 + i * bar_w / 70, bar_y + 14, 7, bar_h - 28, DEEP_TEAL)
    for i in range(30):
        xx = bar_x + 70 + i * (bar_w - 150) / 29
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, emu(xx), emu(bar_y + bar_h + 34), emu(xx), emu(bar_y + bar_h + 145))
        conn.line.color.rgb = rgb(BLUE)
        conn.line.width = line_pt(10)
        add_oval(slide, xx - 15, bar_y + bar_h + 25, 30, 30, BLUE)
    add_oval(slide, cx + 100, density_y + 365, 140, 140, BLUE)
    add_center_text(slide, cx + 100, density_y + 365, 140, 140, "!", 19, WHITE, True)
    add_text(slide, cx + 300, density_y + 400, cw - 420, 90, "Main data challenge: dense greenhouse telemetry, sparse rootzone labels.", 17, "#33657C")


def build_architecture(slide, x, y, w, h):
    cx, cy, cw, _ = section(slide, x, y, w, h, 3, "System Architecture & Workflow")
    add_text(slide, cx, cy, cw, 170, "Support-signal modeling, rootzone feature engineering, and time-aware development.", 25, DARK, True)
    top_y = cy + 300
    box_w, box_h, gap = (cw - 2 * 95) / 3, 420, 95
    top = [
        ("External climate +\nradiation", "temperature - radiation - humidity - ETo", "01"),
        ("Support-signal models", "micro-climate - soil temperature - internal radiation", "02"),
        ("Greenhouse context\nfeatures", "predicted greenhouse context", "03"),
    ]
    for i, (title, desc, n) in enumerate(top):
        bx = cx + i * (box_w + gap)
        add_round(slide, bx, top_y, box_w, box_h, WHITE, BORDER, 3)
        add_text(slide, bx + 55, top_y + 50, box_w - 170, 135, title, 17, TEAL, True)
        add_text(slide, bx + 55, top_y + 190, box_w - 110, 180, desc, 17, TEXT_MUTED)
        add_text(slide, bx + box_w - 150, top_y + 55, 95, 50, n, 13, "#8BBBD0", True, "right")
        if i < 2:
            add_arrow(slide, bx + box_w + 25, top_y + box_h / 2, bx + box_w + gap - 30, top_y + box_h / 2)
    input_y = top_y + box_h + 65
    add_round(slide, cx, input_y, cw, 350, "#F3FBEF", GREEN, 3)
    add_text(slide, cx + 60, input_y + 70, 285, 70, "+ INPUTS", 17, "#3C7A11", True)
    xx = cx + 360
    for label in ["Irrigation & fertigation", "Crop status", "Rootzone state t0"]:
        ww = add_chip(slide, xx, input_y + 58, label, "#4A981A", WHITE, None, 115)
        xx += ww + 45
    add_chip(slide, cx + 70, input_y + 205, "History", "#4A981A", WHITE, None, 105)
    add_center_text(slide, cx + cw / 2 - 675, input_y + 400, 1350, 80, "ROOTZONE FEATURE ENGINEERING", 13, BLUE, True)
    model_y = input_y + 515
    add_round(slide, cx, model_y, cw * 0.72, 500, DEEP_TEAL, DEEP_TEAL, 1)
    add_text(slide, cx + 80, model_y + 75, cw * 0.68 - 120, 95, "Rootzone soft sensor", 22, WHITE, True)
    add_text(slide, cx + 80, model_y + 215, cw * 0.68 - 120, 220, "predicts pH and EC between manual samples - pH/EC anchor, time gap, irrigation history, fertilizer / salt pressure, crop status", 17, WHITE, True)
    add_arrow(slide, cx + cw * 0.72 + 50, model_y + 250, cx + cw * 0.78, model_y + 250, DEEP_TEAL)
    out_x = cx + cw * 0.78
    add_round(slide, out_x, model_y, cx + cw - out_x, 230, GREEN, GREEN, 1)
    add_round(slide, out_x, model_y + 270, cx + cw - out_x, 230, BLUE, BLUE, 1)
    add_center_text(slide, out_x, model_y, cx + cw - out_x, 230, "pH(t)", 19, "#06371A", True)
    add_center_text(slide, out_x, model_y + 270, cx + cw - out_x, 230, "EC(t)", 19, WHITE, True)
    wf_y = y + h - 430
    add_round(slide, cx, wf_y, cw, 295, "#F7FDFF", BORDER, 3)
    add_text(slide, cx + 60, wf_y + 90, 450, 120, "PROJECT\nWORKFLOW", 13, TEAL, True)
    steps = ["Problem\nframing", "Data\ncollection", "Audit\n& EDA", "Master\ndataset", "Support\nmodeling", "Feature\nengineering", "Model\nbuilding", "Validation", "Deployment"]
    sx, sw = cx + 600, (cw - 750) / len(steps)
    for i, s in enumerate(steps):
        fill = "#F1FFF1" if i == len(steps) - 1 else WHITE
        outline = GREEN if i == len(steps) - 1 else BORDER
        bx = sx + i * sw
        add_round(slide, bx, wf_y + 58, sw - 25, 187, fill, outline, 3)
        add_center_text(slide, bx, wf_y + 62, sw - 25, 178, f"{i+1:02d}\n{s}", 13, TEAL, True)


def build_validation(slide, x, y, w, h):
    cx, cy, cw, _ = section(slide, x, y, w, h, 4, "Model Development & Validation")
    add_text(slide, cx, cy, cw, 170, "Trained only on past data, tested on future and skipped samples.", 25, DARK, True)
    y1 = cy + 295
    add_text(slide, cx, y1, cw, 60, "WALK-FORWARD", 13, TEAL, True)
    seg_y, seg_h = y1 + 110, 230
    labels = [("Train 1", True), ("Test 1", False), ("Train 2", True), ("Test 2", False), ("Train 3", True), ("Test 3", False), ("Train 4", True), ("Test 4", False)]
    unit = (cw - 20 * (len(labels) - 1)) / (4 * 1.55 + 4 * 0.75)
    x0 = cx
    for lab, train in labels:
        ww = unit * (1.55 if train else 0.75)
        add_round(slide, x0, seg_y, ww, seg_h, "#CFF4FA" if train else BLUE, "#8CD4E8", 2)
        add_center_text(slide, x0, seg_y, ww, seg_h, lab, 17, TEAL if train else WHITE, True)
        x0 += ww + 20
    y2 = seg_y + 360
    add_text(slide, cx, y2, cw, 60, "SKIPPED-INTERVAL HOLDOUT", 13, TEAL, True)
    bar_y, bw, bh = y2 + 110, cw, 225
    add_round(slide, cx, bar_y, bw, bh, "#DDF7FC", BORDER, 3)
    n, gap = 12, 28
    sw = (bw - gap * (n + 1)) / n
    for i in range(n):
        col = OLIVE if i in [2, 5, 9] else DEEP_TEAL
        add_round(slide, cx + gap + i * (sw + gap), bar_y + 45, sw, bh - 90, col, col, 1)
    add_text(slide, cx, bar_y + bh + 90, cw, 65, "Remove labeled intervals and predict them later.", 17, TEXT_MUTED)
    y3 = bar_y + bh + 285
    add_text(slide, cx, y3, cw, 60, "BASELINE COMPARISON", 13, TEAL, True)
    bx_y, bx_w = y3 + 115, (cw - 80) / 2
    add_round(slide, cx, bx_y, bx_w, 390, WHITE, BORDER, 3)
    add_text(slide, cx + 65, bx_y + 80, bx_w - 120, 70, "Naive carry-forward", 17, DARK, True)
    add_text(slide, cx + 65, bx_y + 205, bx_w - 120, 70, "Last manual pH / EC", 17, DARK)
    add_round(slide, cx + bx_w + 80, bx_y, bx_w, 390, "#F4FCEB", GREEN, 3)
    add_text(slide, cx + bx_w + 145, bx_y + 80, bx_w - 120, 70, "Rootzone soft sensor /", 17, "#1C6C0D", True)
    add_text(slide, cx + bx_w + 145, bx_y + 205, bx_w - 120, 70, "Reduced error vs baseline", 17, "#1C6C0D")
    xx = cx
    for c in ["MAE", "RMSE", "R^2", "gain vs naive", "holdout MAE"]:
        ww = add_chip(slide, xx, bx_y + 500, c, TEAL, "#DDF7FC", None, 135)
        xx += ww + 40


def plot_box(slide, label, img, gx, yy, ww, hh, outline):
    add_text(slide, gx, yy - 70, ww, 55, label, 13, TEXT_MUTED, True)
    add_round(slide, gx, yy, ww, hh, WHITE, outline, 3)
    add_picture_contain(slide, img, gx + 30, yy + 30, ww - 60, hh - 60)


def build_results(slide, x, y, w, h):
    cx, cy, cw, ch = section(slide, x, y, w, h, 5, "Results")
    add_text(slide, cx, cy, cw, 120, "Rootzone soft sensor - actual vs predicted on walk-forward evaluation.", 25, DARK, True)
    ph_index = trim_white(PLOTS / "PH& EC_Results_48H_Version" / "PH" / "ph_actual_vs_pred_index.png")
    ph_scatter = trim_white(PLOTS / "PH& EC_Results_48H_Version" / "PH" / "ph_true_vs_pred.png")
    ec_index = trim_white(PLOTS / "PH& EC_Results_48H_Version" / "EC" / "ec_actual_vs_pred_index.png")
    ec_scatter = trim_white(PLOTS / "PH& EC_Results_48H_Version" / "EC" / "ec_true_vs_pred.png")
    group_gap = 115
    group_w = (cw - group_gap) / 2
    metric_y, metric_h = cy + 265, 575
    groups = [
        (cx, "pH", "PH UNITS - N = 56", TEAL, "#F8FDFF", [("0.330", "MAE", TEAL), ("0.932", "R^2", TEAL), ("62.4%", "GAIN VS\nNAIVE", BLUE)]),
        (cx + group_w + group_gap, "EC", "MS/CM - N = 56", GREEN, "#FAFFF5", [("0.127", "MAE", "#1F5D17"), ("0.982", "R^2", "#1F5D17"), ("35.3%", "GAIN VS\nNAIVE", BLUE)]),
    ]
    for gx, title, units, col, fill, metrics in groups:
        add_round(slide, gx, metric_y, group_w, metric_h, fill, col, 4)
        add_oval(slide, gx + 95, metric_y + 85, 120, 120, col)
        add_text(slide, gx + 255, metric_y + 72, 480, 100, title, 32, col, True)
        add_text(slide, gx + group_w - 1140, metric_y + 115, 1000, 60, units, 13, TEXT_MUTED, True, "right")
        mx_w = (group_w - 320) / 3
        for i, (val, lab, accent) in enumerate(metrics):
            add_metric(slide, gx + 95 + i * (mx_w + 65), metric_y + 250, mx_w, 255, val, lab, accent)
    line_y, line_h = metric_y + metric_h + 175, 880
    scatter_y = line_y + line_h + 205
    scatter_h = ch - (scatter_y - cy) - 85
    plot_box(slide, "ACTUAL VS PREDICTED - INDEX-BASED", ph_index, groups[0][0], line_y, group_w, line_h, TEAL)
    plot_box(slide, "ACTUAL VS PREDICTED - INDEX-BASED", ec_index, groups[1][0], line_y, group_w, line_h, GREEN)
    plot_box(slide, "TRUE VS PREDICTED", ph_scatter, groups[0][0], scatter_y, group_w, scatter_h, TEAL)
    plot_box(slide, "TRUE VS PREDICTED", ec_scatter, groups[1][0], scatter_y, group_w, scatter_h, GREEN)


def build_impact(slide, x, y, w, h):
    cx, cy, cw, _ = section(slide, x, y, w, h, 6, "Conclusions & Impact")
    add_text(slide, cx, cy, cw, 120, "Sparse manual samples are now a usable 10-minute prediction workflow for monitoring and decisions.", 25, DARK, True)
    call_y = cy + 135
    add_round(slide, cx, call_y, cw, 150, "#EAFBFF", "#EAFBFF", 1)
    add_rect(slide, cx, call_y, 28, 150, TEAL)
    add_text(slide, cx + 80, call_y + 38, cw - 130, 90, "Helps growers monitor rootzone conditions continuously and respond earlier to salinity or acidity changes between manual samples.", 20, TEAL)
    impacts = [
        ("Earlier awareness", "Predicts movement between manual samples.", GREEN, "#F4FCEA"),
        ("Decision support", "Supports fertigation, irrigation, and acid dosing.", BLUE, "#EDF5FF"),
        ("Usable workflow", "Packaged as ETL, prediction, and dashboard workflow.", OLIVE, "#FCFCEB"),
    ]
    box_y = call_y + 205
    box_gap = 90
    box_w = (cw - 2 * box_gap) / 3
    for i, (title, desc, col, fill) in enumerate(impacts):
        bx = cx + i * (box_w + box_gap)
        add_round(slide, bx, box_y, box_w, 255, fill, col, 4)
        add_oval(slide, bx + 55, box_y + 58, 100, 100, col)
        add_center_text(slide, bx + 55, box_y + 58, 100, 100, "i", 15, WHITE, True)
        add_text(slide, bx + 190, box_y + 48, box_w - 230, 70, title, 20, "#1B4120", True)
        add_text(slide, bx + 190, box_y + 145, box_w - 235, 85, desc, 17, TEXT_MUTED)


def main():
    POSTER_DIR.mkdir(exist_ok=True)
    prs = Presentation()
    prs.slide_width = Cm(W_CM)
    prs.slide_height = Cm(H_CM)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W, H, "#DCE9EF")
    add_rect(slide, 125, 125, W - 250, H - 250, WHITE)
    ruppin, volcani = load_logos()
    build_header(slide, ruppin, volcani)

    margin = 230
    gap = 120
    y1 = 1180
    row1_h = 2600
    row2_y = 3920
    row2_h = 2600
    row3_y = 6660
    row3_h = 3850
    row4_y = 10650
    row4_h = 1030
    top_w1 = 3300
    top_w2 = W - 2 * margin - gap - top_w1
    build_intro(slide, margin, y1, top_w1, row1_h)
    build_dataset(slide, margin + top_w1 + gap, y1, top_w2, row1_h)
    col_w = (W - 2 * margin - gap) / 2
    build_architecture(slide, margin, row2_y, col_w, row2_h)
    build_validation(slide, margin + col_w + gap, row2_y, col_w, row2_h)
    full_w = W - 2 * margin
    build_results(slide, margin, row3_y, full_w, row3_h)
    build_impact(slide, margin, row4_y, full_w, row4_h)
    prs.save(OUT_PPTX)
    print(OUT_PPTX)
    print(f"{W_CM}x{H_CM} cm, one slide")


if __name__ == "__main__":
    main()
