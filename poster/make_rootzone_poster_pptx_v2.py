from io import BytesIO
from pathlib import Path
from zipfile import ZipFile

import pandas as pd
from PIL import Image, ImageChops, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Cm, Emu, Pt


DPI = 300
W_CM = 70
H_CM = 100
W = round(W_CM / 2.54 * DPI)
H = round(H_CM / 2.54 * DPI)

ROOT = Path(__file__).resolve().parents[1]
POSTER_DIR = ROOT / "poster"
PLOTS = ROOT / "plots"
ASSET_DIR = POSTER_DIR / "assets_v2"
OUT_PPTX = POSTER_DIR / "rootzone_poster_70x100cm_v2_editable.pptx"

BLUE = "#1976D2"
TEAL = "#006B82"
DARK = "#0D2B38"
DEEP_TEAL = "#006276"
GREEN = "#75C856"
OLIVE = "#A6C946"
BORDER = "#9FD4E8"
TEXT_MUTED = "#4D6C7C"
LINE = "#A8D6E8"
WHITE = "#FFFFFF"


def emu(v):
    return Emu(round(v * 914400 / DPI))


def rgb(hex_color):
    return RGBColor.from_string(hex_color.strip("#"))


def line_pt(px):
    return Pt(px * 72 / DPI)


def set_fill(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(color)


def set_line(shape, color=None, width_px=3):
    if color is None:
        shape.line.fill.background()
        return
    shape.line.color.rgb = rgb(color)
    shape.line.width = line_pt(width_px)


def add_round(slide, x, y, w, h, fill, outline=BORDER, width_px=3):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, emu(x), emu(y), emu(w), emu(h))
    set_fill(shp, fill)
    set_line(shp, outline, width_px)
    return shp


def add_rect(slide, x, y, w, h, fill, outline=None, width_px=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, emu(x), emu(y), emu(w), emu(h))
    set_fill(shp, fill)
    set_line(shp, outline, width_px)
    return shp


def add_oval(slide, x, y, w, h, fill, outline=None, width_px=0):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, emu(x), emu(y), emu(w), emu(h))
    set_fill(shp, fill)
    set_line(shp, outline, width_px)
    return shp


def add_text(slide, x, y, w, h, text, size, color=DARK, bold=False, align="left", valign="top", font="Arial"):
    box = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = {
        "top": MSO_ANCHOR.TOP,
        "middle": MSO_ANCHOR.MIDDLE,
        "bottom": MSO_ANCHOR.BOTTOM,
    }.get(valign, MSO_ANCHOR.TOP)
    alignment = {
        "left": PP_ALIGN.LEFT,
        "center": PP_ALIGN.CENTER,
        "right": PP_ALIGN.RIGHT,
    }.get(align, PP_ALIGN.LEFT)
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = alignment
        p.space_after = Pt(0)
        p.space_before = Pt(0)
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def add_center_text(slide, x, y, w, h, text, size, color=DARK, bold=False):
    return add_text(slide, x, y, w, h, text, size, color, bold, "center", "middle")


def add_arrow(slide, x1, y1, x2, y2, color=LINE):
    if x2 <= x1:
        return
    h = 34
    shp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, emu(x1), emu(y1 - h / 2), emu(x2 - x1), emu(h))
    set_fill(shp, color)
    set_line(shp, color, 1)
    return shp


def add_line(slide, x1, y1, x2, y2, color=LINE, width_px=8):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, emu(x1), emu(y1), emu(x2), emu(y2))
    conn.line.color.rgb = rgb(color)
    conn.line.width = line_pt(width_px)
    return conn


def add_down_arrow(slide, x, y1, y2, color=DARK):
    if y2 <= y1:
        return
    w = 54
    shp = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, emu(x - w / 2), emu(y1), emu(w), emu(y2 - y1))
    set_fill(shp, color)
    set_line(shp, color, 1)
    return shp


def add_chevron(slide, x, y, w, h, label, idx, fill=WHITE, outline=BORDER):
    shp = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, emu(x), emu(y), emu(w), emu(h))
    set_fill(shp, fill)
    set_line(shp, outline, 3)
    add_center_text(slide, x + 38, y + 18, w - 72, 44, f"{idx:02d}", 12, TEAL, True)
    add_center_text(slide, x + 42, y + 68, w - 84, h - 86, label, 12, DARK, True)
    return shp


def add_chip(slide, x, y, text, accent=BLUE, fill="#F6FDFF", w=None, h=120):
    if w is None:
        w = max(220, 34 * len(text) + 100)
    add_round(slide, x, y, w, h, fill, accent, 3)
    add_center_text(slide, x + 12, y + 6, w - 24, h - 12, text, 13.5, accent, True)
    return w


def add_metric(slide, x, y, w, h, value, label, accent):
    add_round(slide, x, y, w, h, "#F2FBFE", accent, 4)
    add_center_text(slide, x, y + 12, w, h * 0.56, value, 42, accent, True)
    add_center_text(slide, x + 10, y + h * 0.57, w - 20, h * 0.33, label, 14, TEXT_MUTED, True)


def trim_white(path):
    img = Image.open(path).convert("RGBA")
    rgb_img = img.convert("RGB")
    diff = ImageChops.difference(rgb_img, Image.new("RGB", rgb_img.size, "white"))
    bbox = diff.getbbox()
    if not bbox:
        return img
    pad = 18
    left = max(0, bbox[0] - pad)
    top = max(0, bbox[1] - pad)
    right = min(img.width, bbox[2] + pad)
    bottom = min(img.height, bbox[3] + pad)
    return img.crop((left, top, right, bottom))


def image_bytes(img):
    buf = BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def add_picture_contain(slide, img, x, y, w, h):
    if isinstance(img, (str, Path)):
        im = Image.open(img)
        source = str(Path(img))
    else:
        im = img
        source = image_bytes(img)
    scale = min(w / im.width, h / im.height)
    nw, nh = im.width * scale, im.height * scale
    return slide.shapes.add_picture(source, emu(x + (w - nw) / 2), emu(y + (h - nh) / 2), width=emu(nw), height=emu(nh))


def pil_font(size, bold=False):
    return ImageFont.truetype("arialbd.ttf" if bold else "arial.ttf", size)


def draw_wide_timeseries(csv_path, target, out_path, title, ylabel, accent):
    df = pd.read_csv(csv_path)
    true_col = f"{target}_true"
    pred_col = f"{target}_pred"
    naive_col = f"{target}_naive"
    ts = pd.to_datetime(df["timestamp"])
    true_vals = df[true_col].astype(float).tolist()
    pred_vals = df[pred_col].astype(float).tolist()
    naive_vals = df[naive_col].astype(float).tolist()

    img_w, img_h = 3600, 1060
    ml, mr, mt, mb = 250, 95, 140, 155
    plot_w, plot_h = img_w - ml - mr, img_h - mt - mb
    img = Image.new("RGB", (img_w, img_h), "white")
    d = ImageDraw.Draw(img)
    f_title = pil_font(42, True)
    f_label = pil_font(27, True)
    f_tick = pil_font(22)
    f_leg = pil_font(25, True)

    vals = true_vals + pred_vals + naive_vals
    ymin, ymax = min(vals), max(vals)
    pad = max((ymax - ymin) * 0.12, 0.15 if target == "ec" else 0.4)
    ymin -= pad
    ymax += pad
    if target == "ec":
        ymin = max(0, ymin)

    def sx(i):
        return ml + i * plot_w / (len(true_vals) - 1)

    def sy(v):
        return mt + (ymax - v) * plot_h / (ymax - ymin)

    d.text((ml, 40), title, fill=rgb_to_tuple(DARK), font=f_title)
    for j in range(6):
        y = mt + j * plot_h / 5
        val = ymax - j * (ymax - ymin) / 5
        d.line((ml, y, ml + plot_w, y), fill="#E5F1F6", width=3)
        d.text((55, y - 16), f"{val:.1f}" if target == "ph" else f"{val:.2f}", fill=rgb_to_tuple(TEXT_MUTED), font=f_tick)
    for j in range(0, len(true_vals), max(1, len(true_vals) // 7)):
        x = sx(j)
        d.line((x, mt, x, mt + plot_h), fill="#EEF7FA", width=2)
        d.text((x - 55, mt + plot_h + 36), ts.iloc[j].strftime("%m/%d"), fill=rgb_to_tuple(TEXT_MUTED), font=f_tick)
    d.line((ml, mt, ml, mt + plot_h), fill="#9ACFE0", width=4)
    d.line((ml, mt + plot_h, ml + plot_w, mt + plot_h), fill="#9ACFE0", width=4)

    def series(vals, color, width, dash=False):
        pts = [(sx(i), sy(v)) for i, v in enumerate(vals)]
        if dash:
            for p1, p2 in zip(pts[:-1], pts[1:]):
                x1, y1 = p1
                x2, y2 = p2
                steps = 14
                for k in range(0, steps, 2):
                    xa = x1 + (x2 - x1) * k / steps
                    ya = y1 + (y2 - y1) * k / steps
                    xb = x1 + (x2 - x1) * min(k + 1, steps) / steps
                    yb = y1 + (y2 - y1) * min(k + 1, steps) / steps
                    d.line((xa, ya, xb, yb), fill=color, width=width)
        else:
            d.line(pts, fill=color, width=width, joint="curve")
        for x, y in pts:
            d.ellipse((x - 8, y - 8, x + 8, y + 8), fill=color, outline="white", width=2)

    series(naive_vals, "#A7B7C2", 5, True)
    series(true_vals, "#222222", 5)
    series(pred_vals, accent, 6)
    d.text((16, mt + plot_h / 2 - 20), ylabel, fill=rgb_to_tuple(DARK), font=f_label)
    d.text((ml + plot_w / 2 - 60, img_h - 65), "Measurement date", fill=rgb_to_tuple(DARK), font=f_label)

    leg_x = img_w - 1010
    leg_y = 55
    for i, (name, color, dashed) in enumerate([("Actual", "#222222", False), ("Predicted", accent, False), ("Naive baseline", "#A7B7C2", True)]):
        yy = leg_y + i * 42
        if dashed:
            for k in range(3):
                d.line((leg_x + k * 30, yy + 15, leg_x + k * 30 + 18, yy + 15), fill=color, width=6)
        else:
            d.line((leg_x, yy + 15, leg_x + 75, yy + 15), fill=color, width=6)
        d.text((leg_x + 95, yy), name, fill=rgb_to_tuple(DARK), font=f_leg)

    out_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(out_path, dpi=(DPI, DPI))
    return out_path


def rgb_to_tuple(hex_color):
    h = hex_color.strip("#")
    return tuple(int(h[i : i + 2], 16) for i in (0, 2, 4))


def generate_result_assets():
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    ph_path = draw_wide_timeseries(
        ROOT / "scripts" / "exports" / "v8_final_unified_model_48h_no_rh_eval_ph.csv",
        "ph",
        ASSET_DIR / "v2_ph_wide_timeseries.png",
        "pH actual vs predicted on walk-forward evaluation",
        "pH",
        "#5BC0DE",
    )
    ec_path = draw_wide_timeseries(
        ROOT / "scripts" / "exports" / "v8_final_unified_model_48h_no_rh_eval_ec.csv",
        "ec",
        ASSET_DIR / "v2_ec_wide_timeseries.png",
        "EC actual vs predicted on walk-forward evaluation",
        "EC (mS/cm)",
        "#74C857",
    )
    return ph_path, ec_path


def load_logos():
    ppt = POSTER_DIR / "hypotheticalocean.pptx"
    with ZipFile(ppt) as z:
        logos = []
        for name in z.namelist():
            if name.startswith("ppt/media/"):
                img = Image.open(BytesIO(z.read(name))).convert("RGBA")
                if img.size in [(217, 217), (900, 900)]:
                    logos.append(img)
    ruppin = next(i for i in logos if i.size == (217, 217))
    volcani = next(i for i in logos if i.size == (900, 900))
    return ruppin, volcani


def section(slide, x, y, w, h, n, title):
    add_round(slide, x, y, w, h, WHITE, BORDER, 6)
    add_oval(slide, x + 50, y + 70, 270, 270, BLUE)
    add_center_text(slide, x + 50, y + 70, 270, 270, str(n), 23, WHITE, True)
    add_text(slide, x + 360, y + 88, w - 450, 150, title, 36, TEAL, True)
    return x + 155, y + 345, w - 310, h - 435


def build_header(slide, ruppin, volcani):
    x, y, w, h = 230, 170, W - 460, 900
    add_round(slide, x, y, w, h, WHITE, "#4BA5BF", 7)
    add_picture_contain(slide, ruppin, x + 145, y + 95, 470, 470)
    add_picture_contain(slide, volcani, x + w - 615, y + 90, 470, 470)
    add_center_text(slide, x + 690, y + 80, w - 1380, 360, "Predictive Modeling of Rootzone Variables\nUsing Greenhouse Data", 44, TEAL, True)
    add_center_text(slide, x + 670, y + 485, w - 1340, 110, "Students: Eden Eliyahu, Oren Chaushu", 22, TEAL, True)
    add_center_text(slide, x + 670, y + 635, w - 1340, 110, "Mentors: Dr. Avner Priel - Ruppin,  Dr. Alaa Jamal - Volcani Institute", 18, DARK, True)


def build_intro(slide, x, y, w, h):
    cx, cy, cw, _ = section(slide, x, y, w, h, 1, "Introduction & Goal")
    add_text(slide, cx, cy, cw, 90, "Why rootzone monitoring needs a soft sensor.", 25, DARK, True)
    para = (
        "Rootzone pH and EC are key indicators of nutrient availability, salinity, and plant "
        "health in greenhouse cultivation. Manual samples are accurate, but they are sparse - "
        "between samples, acidity or salinity changes may go undetected."
    )
    add_text(slide, cx, cy + 210, cw, 330, para, 20, TEXT_MUTED)
    box_y = cy + 735
    bh = 790
    bw = (cw - 2 * 190) / 3
    labels = [
        ("Manual\nsample", "accurate -\nsparse", DEEP_TEAL, "#EFFAFF"),
        ("Blind\ninterval", "no observation", "#78910D", "#FBFEEB"),
        ("Rootzone\nsoft-sensor\nprediction", "dense -\nestimated", BLUE, "#F1F7FF"),
    ]
    for i, (title, sub, col, fill) in enumerate(labels):
        bx = cx + i * (bw + 190)
        add_round(slide, bx, box_y, bw, bh, fill, col, 4)
        add_oval(slide, bx + bw / 2 - 82, box_y + 60, 164, 164, fill, col, 16)
        add_center_text(slide, bx + 45, box_y + 280, bw - 90, 260, title, 20, col, True)
        add_center_text(slide, bx + 45, box_y + 535, bw - 90, 180, sub, 17, TEXT_MUTED)
        if i < 2:
            add_arrow(slide, bx + bw + 45, box_y + bh / 2, bx + bw + 125, box_y + bh / 2)
    goal_y = box_y + bh + 120
    add_round(slide, cx, goal_y, cw, 350, "#F4FCFE", BORDER, 4)
    add_oval(slide, cx + 80, goal_y + 135, 65, 65, "#A7E08F")
    add_oval(slide, cx + 96, goal_y + 151, 33, 33, GREEN)
    add_text(slide, cx + 210, goal_y + 75, cw - 300, 210, "Goal - Build a rootzone soft sensor that predicts future pH and EC between manual samples.", 20, TEAL, True)


def build_dataset(slide, x, y, w, h):
    cx, cy, cw, _ = section(slide, x, y, w, h, 2, "Dataset & Processing")
    add_text(slide, cx, cy, cw, 90, "Five greenhouse domains synchronized into one 10-minute master dataset.", 25, DARK, True)
    row_x, row_y = cx, cy + 245
    row_w, row_h, gap = int(cw * 0.55), 210, 42
    rows = [
        ("Climate conditions", "temperature - humidity - radiation", BLUE),
        ("Irrigation & fertigation", "volume - fertilizer - acid & salt inputs", GREEN),
        ("Plant status", "canopy cover - crop age - growth stage", OLIVE),
        ("Rootzone measurements", "manual pH samples - manual EC samples", "#08798B"),
        ("Time & history", "timestamps - gaps - previous state", "#1A5570"),
    ]
    merge_x = row_x + row_w + 330
    merge_y = row_y + 2 * (row_h + gap) + row_h / 2
    for i, (lab, desc, col) in enumerate(rows):
        yy = row_y + i * (row_h + gap)
        add_round(slide, row_x, yy, row_w, row_h, WHITE, BORDER, 3)
        add_oval(slide, row_x + 60, yy + 82, 58, 58, col)
        add_text(slide, row_x + 155, yy + 58, row_w * 0.43, 90, lab, 17, TEAL, True)
        add_text(slide, row_x + row_w * 0.52, yy + 55, row_w * 0.45, 110, desc, 16, TEXT_MUTED)
        add_line(slide, row_x + row_w, yy + row_h / 2, merge_x, merge_y, col, 7)
    add_center_text(slide, merge_x - 85, merge_y - 92, 140, 54, "MERGE", 9, TEAL, True)
    card_x, card_y = cx + int(cw * 0.64), row_y
    card_w, card_h = cw - (card_x - cx), 1220
    add_arrow(slide, merge_x, merge_y, card_x - 28, merge_y, LINE)
    add_round(slide, card_x, card_y, card_w, card_h, DEEP_TEAL, "#034A58", 4)
    add_center_text(slide, card_x, card_y + 80, card_w, 280, "MASTER DATASET\n10-minute\ngreenhouse grid", 20, WHITE, True)
    add_round(slide, card_x + 60, card_y + 410, card_w - 120, card_h - 470, "#197388", "#197388", 1)
    stats = [("16,682", "ROWS"), ("10 min", "GRAIN"), ("109", "PH LABELS"), ("109", "EC LABELS"), ("20", "COLUMNS"), ("May-Sep", "2025")]
    sx = [card_x + 145, card_x + card_w * 0.56]
    sy = card_y + 500
    for i, (val, lab) in enumerate(stats):
        add_text(slide, sx[i % 2], sy + (i // 2) * 275, 430, 80, val, 18, WHITE, True)
        add_text(slide, sx[i % 2], sy + 100 + (i // 2) * 275, 430, 70, lab, 13, "#C6EFF7", True)
    challenge_y = card_y + card_h + 115
    add_round(slide, cx, challenge_y, cw, 360, "#EAFBFF", "#EAFBFF", 1)
    add_oval(slide, cx + 90, challenge_y + 90, 180, 180, BLUE)
    add_center_text(slide, cx + 90, challenge_y + 90, 180, 180, "!", 28, WHITE, True)
    add_text(slide, cx + 330, challenge_y + 72, cw - 430, 130, "Main data challenge:", 26, TEAL, True)
    add_text(slide, cx + 330, challenge_y + 175, cw - 430, 120, "dense greenhouse telemetry, sparse rootzone labels.", 25, "#33657C", True)


def build_architecture(slide, x, y, w, h):
    cx, cy, cw, _ = section(slide, x, y, w, h, 3, "System Architecture & Workflow")
    add_text(slide, cx, cy, cw, 105, "Two views: project workflow and prediction system architecture.", 25, DARK, True)

    workflow_y = cy + 190
    add_text(slide, cx, workflow_y, 780, 65, "PROJECT WORKFLOW", 15, TEAL, True)
    steps = [
        "Problem\nframing", "Data\ncollection", "Audit\n& EDA", "Master\ndataset",
        "Support\nmodeling", "Feature\neng.", "Model\nbuilding", "Validation", "Deployment"
    ]
    chevron_y = workflow_y + 78
    gap = 18
    sw = (cw - gap * 8) / 9
    for i, s in enumerate(steps):
        fill = "#F1FFF1" if i == 8 else WHITE
        outline = GREEN if i == 8 else BORDER
        add_chevron(slide, cx + i * (sw + gap), chevron_y, sw, 260, s, i + 1, fill, outline)

    arch_y = chevron_y + 390
    add_text(slide, cx, arch_y, 1050, 65, "SYSTEM ARCHITECTURE", 15, TEAL, True)
    canvas_y = arch_y + 78
    add_round(slide, cx, canvas_y, cw, 1685, "#EEF9FC", "#D6EEF6", 2)
    left_x = cx + 80
    mid_x = cx + cw * 0.42
    right_x = cx + cw * 0.72
    box_h = 300
    input_w = cw * 0.28
    input_rows = [
        ("External Weather", "Bet Dagan: temp, radiation, humidity", "#D8EAFB", BLUE, canvas_y + 115),
        ("Irrigation Plan", "volume, fertilizer concentration, recency", "#D8EAFB", BLUE, canvas_y + 600),
        ("Current State", "last measured EC / pH @ t0", "#D8EAFB", BLUE, canvas_y + 1085),
    ]
    for title, desc, fill, outline, yy in input_rows:
        add_round(slide, left_x, yy, input_w, box_h, fill, outline, 4)
        add_center_text(slide, left_x + 45, yy + 82, input_w - 90, 75, title, 18, DARK, True)
        add_center_text(slide, left_x + 45, yy + 185, input_w - 90, 60, desc, 11, TEXT_MUTED)

    micro_w = cw * 0.28
    pred_w = cw * 0.28
    add_round(slide, mid_x, canvas_y + 115, micro_w, box_h, "#DDF4E3", GREEN, 4)
    add_center_text(slide, mid_x + 45, canvas_y + 205, micro_w - 90, 75, "Micro-Climate Model", 18, DARK, True)
    add_center_text(slide, mid_x + 45, canvas_y + 300, micro_w - 90, 60, "LightGBM", 11, TEXT_MUTED)
    add_round(slide, right_x, canvas_y + 115, pred_w, box_h, "#FFF0CE", "#E8A92B", 4)
    add_center_text(slide, right_x + 45, canvas_y + 205, pred_w - 90, 75, "Predicted Internal Climate", 18, DARK, True)
    add_center_text(slide, right_x + 45, canvas_y + 300, pred_w - 90, 60, "temperature and ET0", 11, TEXT_MUTED)
    root_y = canvas_y + 650
    add_round(slide, right_x, root_y, pred_w, 350, "#D7F0EC", TEAL, 4)
    add_center_text(slide, right_x + 45, root_y + 95, pred_w - 90, 75, "Rootzone Model", 18, DARK, True)
    add_center_text(slide, right_x + 45, root_y + 205, pred_w - 90, 70, "hybrid linear + boosted trees", 11, TEXT_MUTED)
    pred_y = canvas_y + 1165
    add_round(slide, right_x, pred_y, pred_w, 350, "#FCE0D9", "#ED765E", 4)
    add_center_text(slide, right_x + 45, pred_y + 95, pred_w - 90, 75, "Prediction", 18, DARK, True)
    add_center_text(slide, right_x + 45, pred_y + 205, pred_w - 90, 70, "EC & pH @ t+1", 11, TEXT_MUTED)

    add_arrow(slide, left_x + input_w, canvas_y + 265, mid_x - 35, canvas_y + 265, DARK)
    add_arrow(slide, mid_x + micro_w, canvas_y + 265, right_x - 35, canvas_y + 265, DARK)
    add_down_arrow(slide, right_x + pred_w / 2, canvas_y + 415, root_y - 10, DARK)
    add_line(slide, left_x + input_w, canvas_y + 750, right_x - 35, root_y + 130, DARK, 9)
    add_arrow(slide, right_x - 120, root_y + 130, right_x - 35, root_y + 130, DARK)
    add_line(slide, left_x + input_w, canvas_y + 1235, right_x - 35, root_y + 250, DARK, 9)
    add_arrow(slide, right_x - 120, root_y + 250, right_x - 35, root_y + 250, DARK)
    add_down_arrow(slide, right_x + pred_w / 2, root_y + 350, pred_y - 10, DARK)


def build_validation(slide, x, y, w, h):
    cx, cy, cw, _ = section(slide, x, y, w, h, 4, "Model Development & Validation")
    add_text(slide, cx, cy, cw, 145, "Validation separates future testing, skipped intervals, and hybrid model behavior.", 22, DARK, True)

    wf_y = cy + 245
    add_text(slide, cx, wf_y, cw, 55, "WALK-FORWARD STAIRCASE", 13, TEAL, True)
    step_h, test_w = 118, 230
    train_unit = (cw - test_w - 70) / 4
    for i in range(4):
        yy = wf_y + 82 + i * (step_h + 28)
        train_w = train_unit * (i + 1)
        add_round(slide, cx, yy, train_w, step_h, "#CFF4FA", "#8CD4E8", 2)
        add_center_text(slide, cx, yy, train_w, step_h, f"Train 1-{i+1}" if i else "Train 1", 14, TEAL, True)
        add_round(slide, cx + train_w + 35, yy, test_w, step_h, BLUE, BLUE, 2)
        add_center_text(slide, cx + train_w + 35, yy, test_w, step_h, f"Test {i+1}", 14, WHITE, True)

    hold_y = wf_y + 82 + 4 * (step_h + 28) + 45
    add_text(slide, cx, hold_y, cw, 55, "SKIPPED-INTERVAL HOLDOUT", 13, TEAL, True)
    bar_y, bh = hold_y + 80, 190
    add_round(slide, cx, bar_y, cw, bh, "#DDF7FC", BORDER, 3)
    n, gap = 10, 24
    sw = (cw - gap * (n + 1)) / n
    for i in range(n):
        col = OLIVE if i in [2, 6, 8] else DEEP_TEAL
        add_round(slide, cx + gap + i * (sw + gap), bar_y + 45, sw, bh - 90, col, col, 1)
    add_text(slide, cx, bar_y + bh + 55, cw, 60, "Hold out labeled intervals, then predict them later.", 15, TEXT_MUTED)

    hybrid_y = bar_y + bh + 175
    add_text(slide, cx, hybrid_y, cw, 55, "ROOTZONE HYBRID MODEL", 13, TEAL, True)
    card_y = hybrid_y + 75
    card_w = (cw - 120) / 2
    add_round(slide, cx, card_y, card_w, 255, "#F7FDFF", BORDER, 3)
    add_text(slide, cx + 45, card_y + 50, card_w - 90, 60, "Robust linear model", 16, TEAL, True)
    add_text(slide, cx + 45, card_y + 125, card_w - 90, 90, "Interpolates local pH/EC movement between nearby manual samples.", 13, TEXT_MUTED)
    add_round(slide, cx + card_w + 120, card_y, card_w, 255, "#F4FCEB", GREEN, 3)
    add_text(slide, cx + card_w + 165, card_y + 50, card_w - 90, 60, "Gradient boosted trees", 16, "#1C6C0D", True)
    add_text(slide, cx + card_w + 165, card_y + 125, card_w - 90, 90, "Handles extrapolation with climate, irrigation, crop and time features.", 13, TEXT_MUTED)
    merge_y = card_y + 365
    add_line(slide, cx + card_w / 2, card_y + 255, cx + cw / 2, merge_y, TEAL, 6)
    add_line(slide, cx + card_w + 120 + card_w / 2, card_y + 255, cx + cw / 2, merge_y, GREEN, 6)
    add_round(slide, cx + cw * 0.18, merge_y, cw * 0.64, 165, DEEP_TEAL, DEEP_TEAL, 2)
    add_center_text(slide, cx + cw * 0.18, merge_y, cw * 0.64, 165, "Hybrid rootzone soft sensor\npH(t) + EC(t)", 16, WHITE, True)


def plot_box(slide, label, img, gx, yy, ww, hh, outline):
    add_text(slide, gx, yy - 70, ww, 55, label, 13, TEXT_MUTED, True)
    add_round(slide, gx, yy, ww, hh, WHITE, outline, 3)
    add_picture_contain(slide, img, gx + 30, yy + 30, ww - 60, hh - 60)


def build_results(slide, x, y, w, h):
    cx, cy, cw, ch = section(slide, x, y, w, h, 5, "Results")
    add_text(slide, cx, cy, cw, 115, "Wide walk-forward traces from export data: actual, predicted, and naive baseline.", 24, DARK, True)
    ph_plot, ec_plot = generate_result_assets()
    metric_y = cy + 205
    block_gap = 90
    block_w = (cw - block_gap) / 2
    metrics = [
        (cx, "pH", TEAL, "#F8FDFF", [("0.330", "MAE", TEAL), ("0.932", "R^2", TEAL), ("62.4%", "GAIN", BLUE)]),
        (cx + block_w + block_gap, "EC", GREEN, "#FAFFF5", [("0.127", "MAE", "#1F5D17"), ("0.982", "R^2", "#1F5D17"), ("35.3%", "GAIN", BLUE)]),
    ]
    for gx, title, col, fill, vals in metrics:
        add_round(slide, gx, metric_y, block_w, 445, fill, col, 4)
        add_oval(slide, gx + 80, metric_y + 80, 120, 120, col)
        add_text(slide, gx + 235, metric_y + 72, 320, 100, title, 31, col, True)
        mw = (block_w - 315) / 3
        for i, (val, lab, accent) in enumerate(vals):
            add_metric(slide, gx + 80 + i * (mw + 58), metric_y + 210, mw, 185, val, lab, accent)

    plot_y = metric_y + 560
    plot_gap = 170
    plot_h = (ch - (plot_y - cy) - plot_gap - 55) / 2
    plot_box(slide, "pH ACTUAL VS PREDICTED - WALK-FORWARD", ph_plot, cx, plot_y, cw, plot_h, TEAL)
    plot_box(slide, "EC ACTUAL VS PREDICTED - WALK-FORWARD", ec_plot, cx, plot_y + plot_h + plot_gap, cw, plot_h, GREEN)


def build_impact(slide, x, y, w, h):
    cx, cy, cw, _ = section(slide, x, y, w, h, 6, "Conclusions & Impact")
    add_text(slide, cx, cy, cw, 190, "Sparse manual measurements transformed into a usable prediction workflow.", 23, DARK, True)
    bullets = [
        ("Unified 10-minute master dataset", "multiple greenhouse domains"),
        ("Support-signal modeling", "completed greenhouse context"),
        ("Hybrid rootzone soft sensor", "reduced error vs baseline"),
        ("Packaged workflow", "ETL, prediction and dashboard"),
    ]
    yy = cy + 245
    for title, desc in bullets:
        add_oval(slide, cx + 5, yy + 12, 58, 58, BLUE)
        add_oval(slide, cx + 20, yy + 27, 28, 28, "#BFE8FF")
        add_text(slide, cx + 88, yy, cw - 105, 62, title, 15, DARK, True)
        add_text(slide, cx + 88, yy + 62, cw - 105, 54, desc, 13, TEXT_MUTED)
        yy += 170

    call_y = yy + 20
    add_round(slide, cx, call_y, cw, 360, "#EAFBFF", "#EAFBFF", 1)
    add_rect(slide, cx, call_y, 28, 360, TEAL)
    add_text(slide, cx + 78, call_y + 58, cw - 120, 230, "Helps growers monitor rootzone conditions continuously and respond earlier to salinity or acidity changes between manual samples.", 18, TEAL)

    impacts = [
        ("Earlier awareness", "pH and EC movement between samples", GREEN, "#F4FCEA"),
        ("Decision support", "fertigation, irrigation and acid dosing", BLUE, "#EDF5FF"),
        ("Usable workflow", "automated ETL and dashboard packaging", OLIVE, "#FCFCEB"),
    ]
    box_y = call_y + 455
    for title, desc, col, fill in impacts:
        add_round(slide, cx, box_y, cw, 315, fill, col, 4)
        add_oval(slide, cx + 62, box_y + 78, 105, 105, col)
        add_center_text(slide, cx + 62, box_y + 78, 105, 105, "i", 16, WHITE, True)
        add_text(slide, cx + 205, box_y + 55, cw - 245, 72, title, 18, "#1B4120", True)
        add_text(slide, cx + 205, box_y + 150, cw - 245, 100, desc, 15, TEXT_MUTED)
        box_y += 385


def main():
    POSTER_DIR.mkdir(exist_ok=True)
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = Cm(W_CM)
    prs.slide_height = Cm(H_CM)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_rect(slide, 0, 0, W - 1, H, "#DCE9EF")
    add_rect(slide, 125, 125, W - 250, H - 250, WHITE)
    ruppin, volcani = load_logos()
    build_header(slide, ruppin, volcani)

    margin = 230
    gap = 120
    y1 = 1180
    row1_h = 2500
    row2_y = 3820
    row2_h = 3050
    row3_y = 7010
    row3_h = 4570
    top_w1 = 3300
    top_w2 = W - 2 * margin - gap - top_w1
    build_intro(slide, margin, y1, top_w1, row1_h)
    build_dataset(slide, margin + top_w1 + gap, y1, top_w2, row1_h)
    arch_w = 5200
    validation_w = W - 2 * margin - gap - arch_w
    build_architecture(slide, margin, row2_y, arch_w, row2_h)
    build_validation(slide, margin + arch_w + gap, row2_y, validation_w, row2_h)
    impact_w = 1900
    results_w = W - 2 * margin - gap - impact_w
    build_results(slide, margin, row3_y, results_w, row3_h)
    build_impact(slide, margin + results_w + gap, row3_y, impact_w, row3_h)
    prs.save(OUT_PPTX)
    print(OUT_PPTX)
    print(f"{W_CM}x{H_CM} cm, one slide")


if __name__ == "__main__":
    main()
