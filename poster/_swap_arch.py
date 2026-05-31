"""Swap ONLY the Method architecture picture (#36) in the Project 14 deliverable
for the regenerated big-text architecture.png, at the same position/size. Touches
nothing else.  Run:  py -X utf8 poster/_swap_arch.py
"""
import os
import shutil
from PIL import Image
from pptx import Presentation
from pptx.util import Cm

ROOT  = r"C:\Users\User\Desktop\Predictive_Modeling_Of_Rootzone_Variables_Using_Greenhouse_Data"
DELIV = os.path.join(ROOT, "poster",
                     "Predictive Modeling of Root-Zone Variables Using Greenhouse Data - Project 14.pptx")
BAK   = os.path.join(ROOT, "poster",
                     "Predictive Modeling of Root-Zone Variables Using Greenhouse Data - Project 14.PRE_ARCH3.bak.pptx")
ARCH  = os.path.join(ROOT, "poster", "_figs", "architecture.png")

shutil.copy2(DELIV, BAK)
print("backup ->", BAK)

prs = Presentation(DELIV)
slide = prs.slides[0]
byid = {sh.shape_id: sh for sh in slide.shapes}

old = byid[36]
L, T, W = old.left, old.top, old.width          # keep exact position + width
old._element.getparent().remove(old._element)

iw, ih = Image.open(ARCH).size
W_cm = round(W / 360000, 2)
H_cm = round(W_cm * ih / iw, 2)
slide.shapes.add_picture(ARCH, L, T, W, Cm(H_cm))
print(f"architecture: x={round(L/360000,2)} y={round(T/360000,2)} w={W_cm} h={H_cm}")

prs.save(DELIV)
print("saved", DELIV)
