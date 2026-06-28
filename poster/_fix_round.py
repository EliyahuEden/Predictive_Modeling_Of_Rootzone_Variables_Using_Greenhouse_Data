# -*- coding: utf-8 -*-
"""Make all 10 section cards have the SAME visual corner radius.
roundRect radius = min(w,h) * adj/100000. Cards share width 32.4cm but vary in
height (= the shorter side), so equal adj gave unequal radii. Fix: per-card adj
so radius == TARGET for every card. Only the 10 section cards are touched.
"""
from pptx import Presentation
ns = {"a":"http://schemas.openxmlformats.org/drawingml/2006/main"}
A = "http://schemas.openxmlformats.org/drawingml/2006/main"

F = "Predictive_Modeling_Root-Zone_Poster_Final.pptx"
TARGET_CM = 0.6
TARGET_EMU = TARGET_CM * 360000.0

prs = Presentation(F)
slide = prs.slides[0]
changed = []
for sh in slide.shapes:
    el = sh._element
    geom = el.find(".//a:prstGeom", ns)
    if geom is None or geom.get("prst") != "roundRect":
        continue
    w, h = sh.width, sh.height
    # section cards: full-column width ~32.4cm
    if abs(w/360000.0 - 32.4) > 0.15:
        continue
    ss = min(w, h)
    adj = round(TARGET_EMU / ss * 100000.0)
    adj = max(0, min(50000, adj))
    gd = geom.find("a:avLst/a:gd", ns)
    if gd is None:
        av = geom.find("a:avLst", ns)
        if av is None:
            from lxml import etree
            av = etree.SubElement(geom, f"{{{A}}}avLst")
        from lxml import etree
        gd = etree.SubElement(av, f"{{{A}}}gd"); gd.set("name","adj")
    old = gd.get("fmla")
    gd.set("fmla", f"val {adj}")
    radius_cm = ss/360000.0 * adj/100000.0
    changed.append((sh.shape_id, round(h/360000.0,2), old, f"val {adj}", round(radius_cm,3)))

prs.save(F)
print(f"TARGET radius = {TARGET_CM} cm | cards changed: {len(changed)}")
print(f"{'id':>4} {'H_cm':>6} {'old':>10} {'new':>10} {'radius_cm':>9}")
for c in sorted(changed):
    print(f"{c[0]:>4} {c[1]:>6} {c[2]:>10} {c[3]:>10} {c[4]:>9}")
