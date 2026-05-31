"""Surgically rebuild ONLY the Root-zone Model Results card (08) on the Project 14 base.
Reorganises into: walk-forward (intro + compact metric chips + 2 index time-series)
then holdout (intro + 2 parity scatters). Touches nothing outside card 08.
Run:  py -X utf8 poster/_build_results.py
"""
import os
import shutil
from PIL import Image
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR

ROOT  = r"C:\Users\User\Desktop\Predictive_Modeling_Of_Rootzone_Variables_Using_Greenhouse_Data"
DELIV = os.path.join(ROOT, "poster",
                     "Predictive Modeling of Root-Zone Variables Using Greenhouse Data - Project 14.pptx")
BAK   = os.path.join(ROOT, "poster",
                     "Predictive Modeling of Root-Zone Variables Using Greenhouse Data - Project 14.PRE_RESULTS.bak.pptx")
FIGS  = os.path.join(ROOT, "poster", "_figs")

shutil.copy2(DELIV, BAK)
print("backup ->", BAK)

prs = Presentation(DELIV)
slide = prs.slides[0]
byid = {sh.shape_id: sh for sh in slide.shapes}

CX = 36.65          # left text x inside card
CW = 30.7           # inner content width
PW = 15.2           # plot width
XL, XR = 36.55, 52.25   # left / right plot x


def add_intro(text, top, h=2.5):
    tb = slide.shapes.add_textbox(Cm(CX), Cm(top), Cm(CW), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2.8); tf.margin_right = Pt(2.8)
    tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.line_spacing = 1.06
    r = p.add_run(); r.text = text
    r.font.name = "Calibri"; r.font.size = Pt(24)
    r.font.bold = False; r.font.italic = True
    r.font.color.theme_color = MSO_THEME_COLOR.TEXT_2
    return tb


def place(fname, x, top):
    p = os.path.join(FIGS, fname)
    iw, ih = Image.open(p).size
    h = round(PW * ih / iw, 2)
    slide.shapes.add_picture(p, Cm(x), Cm(top), Cm(PW), Cm(h))
    print(f"  {fname}: x={x} y={top} w={PW} h={h} (bottom {round(top+h,2)})")
    return h


# ---------- 1. walk-forward intro sentence ----------
add_intro("Walk-forward validation shows that the model predicts future pH and EC "
          "from past observations while maintaining strong accuracy over time.", 36.1)

# ---------- 2. compact metric chips: move existing #54-59 down ----------
for cid in (54, 55, 56):
    byid[cid].top = Cm(39.0)
for cid in (57, 58, 59):
    byid[cid].top = Cm(41.3)
print("chips repositioned: row1 T=39.0, row2 T=41.3 (bottom 43.4)")

# ---------- 3. remove old parity scatters #60 / #61 ----------
for pid in (60, 61):
    el = byid[pid]._element
    el.getparent().remove(el)
print("removed old parity scatters #60, #61")

# ---------- 4. walk-forward index time-series (replace scatters) ----------
print("walk-forward time-series:")
hts = place("wf_ph_index.png", XL, 44.0)
place("wf_ec_index.png", XR, 44.0)

# ---------- 5. holdout intro sentence ----------
add_intro("Skipped-interval holdout testing confirms that the model generalizes to "
          "unseen intervals that were removed from the training process.", 55.0)

# ---------- 6. holdout parity scatters ----------
print("holdout scatters:")
place("holdout_ph_parity.png", XL, 58.0)
place("holdout_ec_parity.png", XR, 58.0)

prs.save(DELIV)
print("saved", DELIV)
